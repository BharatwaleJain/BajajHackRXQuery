import os
import fitz
import docx
from langchain_community.document_loaders import UnstructuredEmailLoader
import asyncio
import io
import logging
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from bs4 import BeautifulSoup
from app.core.config import settings
def _extract_text_from_html_sync(html_content: str) -> str:
    soup = BeautifulSoup(html_content, "html.parser")
    return soup.get_text(separator="\n", strip=True)
def _extract_text_from_pptx_sync(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    text = []
    tesseract_cmd = settings.tesseract_cmd
    ocr_enabled = bool(tesseract_cmd and os.path.exists(tesseract_cmd))
    if ocr_enabled:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
    else:
        logging.warning("OCR skipped for PPTX images: TESSERACT_CMD not set or path not found.")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
            if shape.shape_type == 13 and ocr_enabled:
                image = shape.image
                image_bytes = image.blob
                try:
                    img = Image.open(io.BytesIO(image_bytes)).convert('L')
                    ocr_text = pytesseract.image_to_string(img, config='--psm 6')
                    if ocr_text:
                        text.append(ocr_text)
                except Exception as e:
                    logging.error(f"Error processing image in PPTX: {e}")
    return "\n".join(text)
def _extract_text_from_xlsx_sync(xlsx_path: str) -> str:
    workbook = load_workbook(filename=xlsx_path)
    text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text.append(str(cell))
    return "\n".join(text)
def _extract_text_from_image_sync(image_path: str) -> str:
    tesseract_cmd = settings.tesseract_cmd
    if not tesseract_cmd or not os.path.exists(tesseract_cmd):
        logging.warning(
            "OCR skipped: Tesseract not found at TESSERACT_CMD. "
            "Set TESSERACT_CMD in environment to the tesseract executable path."
        )
        return ""
    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
    try:
        img = Image.open(image_path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logging.error(f"OCR error on image '{image_path}': {e}")
        return ""
def _extract_text_from_pdf_sync(pdf_path: str) -> str:
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text
def _extract_text_from_docx_sync(docx_path: str) -> str:
    doc = docx.Document(docx_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])
async def extract_text(file_path: str) -> str:
    _, extension = os.path.splitext(file_path)
    ext_lower = extension.lower()
    if ext_lower == ".pdf":
        return await extract_text_from_pdf(file_path)
    elif ext_lower == ".docx":
        return await extract_text_from_docx(file_path)
    elif ext_lower == ".eml":
        return await extract_text_from_eml(file_path)
    elif ext_lower == ".pptx":
        return await extract_text_from_pptx(file_path)
    elif ext_lower == ".xlsx":
        return await extract_text_from_xlsx(file_path)
    elif ext_lower in [".png", ".jpeg", ".jpg"]:
        return await extract_text_from_image(file_path)
    elif ext_lower in [".zip", ".bin"]:
        raise ValueError(f"Unsupported and skipped file type: {extension}")
    elif ext_lower in [".html", ".htm"]:
        with open(file_path, "r", encoding="utf-8") as f:
            html_content = f.read()
        return await extract_text_from_html(html_content)
    else:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            return await extract_text_from_html(content)
        except (UnicodeDecodeError, IsADirectoryError):
             raise ValueError(f"Unsupported file type: {extension}")
async def extract_text_from_pptx(pptx_path: str) -> str:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _extract_text_from_pptx_sync, pptx_path)
async def extract_text_from_xlsx(xlsx_path: str) -> str:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _extract_text_from_xlsx_sync, xlsx_path)
async def extract_text_from_image(image_path: str) -> str:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _extract_text_from_image_sync, image_path)
async def extract_text_from_pdf(pdf_path: str) -> str:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _extract_text_from_pdf_sync, pdf_path)
async def extract_text_from_docx(docx_path: str) -> str:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _extract_text_from_docx_sync, docx_path)
async def extract_text_from_eml(eml_path: str) -> str:
    loader = UnstructuredEmailLoader(eml_path)
    data = await loader.aload()
    return data[0].page_content
async def extract_text_from_html(html_content: str) -> str:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _extract_text_from_html_sync, html_content)